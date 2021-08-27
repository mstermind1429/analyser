import base64
import csv
import logging
import requests
from json import JSONDecodeError
from operator import itemgetter
from collections import Counter

from django.conf import settings
from django.db import IntegrityError

from background_task import background
from lighthouse import LighthouseRunner
import pandas as pd
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.chart import XL_DATA_LABEL_POSITION
from pptx.util import Pt
from pptx.dml.color import RGBColor
from urllib.parse import urlparse

from core.helpers import (generate_random_number, generate_directories, get_keyword_volume, get_subdomain, average,
                          get_access_token, log_to_telegram_bot, send_mail)
from core.searchmetrics import SearchmetricsAPI
from .models import *

logger = logging.getLogger('django')

numeric_aspects = {"first-contentful-paint": "first_cp", "largest-contentful-paint": "largest_cp",
                   "first-meaningful-paint": "first_mp", "speed-index": "speed_index",
                   "estimated-input-latency": "estimated_il", "total-blocking-time": "total_blocking_time",
                   "max-potential-fid": "max_potential_fid", "cumulative-layout-shift": "cumulative_ls",
                   "server-response-time": "server_response_time", "first-cpu-idle": "first_cpu_idle",
                   "interactive": "interactive", "redirects": "redirects",
                   "mainthread-work-breakdown": "mainthread_work_breakdown", "bootup-time": "bootup_time",
                   "uses-rel-preload": "uses_rel_preload", "network-rtt": "network_rtt",
                   "network-server-latency": "network_sl", "uses-long-cache-ttl": "long_cache_ttl"}

score_aspects = {
    "largest-contentful-paint": "lcp_score",
    "max-potential-fid": "fid_score",
    "cumulative-layout-shift": "cls_score"
}

items = {"bootup-time": LighthouseReportBootupItems, "diagnostics": LighthouseReportDiagnosticsItems,
         "network-requests": LighthouseReportNetworkReqItems, "network-rtt": LighthouseReportNetworkRttItems,
         "network-server-latency": LighthouseReportNetworkSLItems, "long-tasks": LighthouseReportLongTaskItem,
         "offscreen-images": LighthouseReportOffscreenImagesItem, "unminified-javascript": LighthouseReportUnminifiedJS,
         "unminified-css": LighthouseReportUnminifiedCSS, "unused-css-rules": LighthouseReportUnusedCSSRule}

FILE_PATH = settings.REPORT_PATH + '/urls_lighthouse.xlsx'


def lighthouse_report(url):
    report_types = ['mobile', 'desktop']
    domain = urlparse(url).netloc
    report_name = f'lighthouse_{domain}_{generate_random_number()}_'

    serializer = {'desktop_report_path': settings.REPORT_URL + report_name + 'desktop.json',
                  'mobile_report_path': settings.REPORT_URL + report_name + 'mobile.json',
                  'data': []}

    for report_type in report_types:
        report_path = f'{settings.REPORT_PATH}/{report_name}{report_type}.json'

        report = LighthouseRunner(url, form_factor=report_type, quiet=False,
                                  report_path=report_path).report

        try:
            snippet = report.get_audits()['largest-contentful-paint-element']['details']['items'][0]['node']['snippet']
            largest_cp_element = snippet.split("src=\"")[1].split("\"")[0]
            content_type = 'image'
        except Exception as e:
            logger.info(f"No image found for the LCP. Error: {e}")
            node = report.get_audits()['largest-contentful-paint-element']['details']['items'][0]['node']['nodeLabel']
            largest_cp_element = node
            content_type = 'node'

        data = {'paint_score': report.get_audits()['largest-contentful-paint']['numericValue'],
                'cumulative_ls': report.get_audits()['cumulative-layout-shift']['numericValue'],
                'max_potential_fid': report.get_audits()['max-potential-fid']['numericValue'],
                'largest_cp_element': largest_cp_element, 'content_type': content_type,
                'audits': report.audits}

        file_name = f'{report_name}{report_type}.csv'
        with open(f'{settings.REPORT_PATH}/{file_name}', 'w+', newline='') as f:
            writer = csv.writer(f)
            writer.writerow([domain, report_type])

            for key in data:
                if key is not 'audits':
                    writer.writerow([key, data[key]])
                else:
                    for audit in data['audits']['performance']:
                        for aspect in audit:
                            writer.writerow([aspect[0], aspect[1], aspect[2]])

        data['csv'] = file_name
        serializer['data'].append(data)

    return serializer


@background(schedule=1)
def run_domain_lighthouse(key, secret, amount, country_code, domain, uploaded_file_url=None):
    api = SearchmetricsAPI(key, secret)

    if not uploaded_file_url:
        keyword_info = {}
        keywords = []
        offset = 0
        for _ in range(amount // 250):
            status, response = api.get_rankings_domain(
                domain=domain,
                country_code=country_code,
                offset=offset
            )

            for element in response:
                keywords.append(element['keyword'])
                keyword_info[element['keyword']] = element

            offset += 250

        ranking = []
        keywords_list, position, urls = [], [], []
        for index, keyword in enumerate(keywords):
            status, ranking_info = api.get_list_rankings_keyword(
                keyword=keyword,
                country_code=country_code
            )

            if not status:
                continue

            if len(ranking_info) == 0:
                ranking.append(keyword)

            for info in ranking_info:
                keywords_list.append(keyword)
                position.append(info['position'])
                urls.append(info['url'])

        df = pd.DataFrame(
            {
                'Keyword': keywords_list,
                'URL': urls,
                'Position': position,
            }
        )
    else:
        df = pd.read_excel(f'/{uploaded_file_url}')
        keywords_list = df['Keyword']
        amount = len(keywords_list)

    if len(keywords_list) == 0:
        send_mail('Domain Lighthouse', f'No data received for domain {domain}. Values set:\n'
                                       f'Country code: {country_code}\nAmount: {amount}')
        return

    run_lighthouse_report(df, domain, amount)


def run_lighthouse_report(df, domain, amount):
    failed, completed, already_exists = [], [], []
    for url, keyword, position in zip(df["URL"], df["Keyword"], df["Position"]):
        url = f"https://{url}"
        report_types = ['mobile']

        for report_type in report_types:
            try:
                report = LighthouseRunner(url, form_factor=report_type, quiet=False).report
            except AttributeError:
                failed.append(url)
                continue

            audits = report.get_audits()
            content_type, node = None, None
            try:
                node = audits['largest-contentful-paint-element']['details']['items'][0]['node']['snippet']
                node.replace('{', '').replace('}', '')
                content_type = 'image'
            except (IndexError, KeyError, TypeError) as err:
                logger.info(f'Error occurred while getting LCP: {err}')

            try:
                node = audits['largest-contentful-paint-element']['details']['items'][0]['node']['nodeLabel']
                node.replace('{', '').replace('}', '')
                content_type = 'node'
            except (IndexError, KeyError, TypeError) as err:
                logger.info(f'Error occurred while getting LCP: {err}')

            lighthouse_report = {}
            for key in numeric_aspects:
                try:
                    lighthouse_report[numeric_aspects[key]] = round(audits[key]['numericValue'], 3)
                except KeyError:
                    failed.append(url)
                    continue

            for key in score_aspects:
                lighthouse_report[score_aspects[key]] = audits[key]['score']

            lighthouse_report['domain'] = domain
            lighthouse_report['url'] = url
            lighthouse_report['keyword'] = keyword
            lighthouse_report['type'] = report_type
            lighthouse_report['position'] = position
            lighthouse_report['content_type'] = content_type
            lighthouse_report['largest_cp_element'] = node

            try:
                lighthouse_report = LighthouseReport.objects.create(**lighthouse_report)
            except IntegrityError:
                lighthouse_report = LighthouseReport.objects.filter(domain='taz.de', url=url, keyword=keyword)

                if not lighthouse_report:
                    failed.append(url)
                else:
                    already_exists.append(url)
                continue

            completed.append(url)

            for key in ["screenshot-thumbnails", "final-screenshot"]:
                try:
                    for item in audits[key]['details']['items']:
                        lighthouse_item = {}
                        for aspect in item:
                            if type(item[aspect]) == str:
                                if '{' in item[aspect] or '}' in item[aspect]:
                                    continue
                            if aspect == 'data':
                                lighthouse_item["content_type"] = item[aspect].split(";")[0]
                                lighthouse_item["data"] = item[aspect].split(";")[1]
                                continue
                            lighthouse_item[aspect] = item[aspect]

                        lighthouse_item["report"] = lighthouse_report

                        LighthouseReportItem.objects.create(**lighthouse_item)
                except (KeyError, ValueError):
                    continue

            for key in items:
                try:
                    for item in audits[key]['details']['items']:
                        lighthouse_item = {}
                        for aspect in item:
                            if type(item[aspect]) == str:
                                if '{' in item[aspect] or '}' in item[aspect]:
                                    continue
                            lighthouse_item[aspect] = round(item[aspect], 2) if type(item[aspect]) == float \
                                else item[aspect]
                        lighthouse_item["report"] = lighthouse_report

                        items[key].objects.create(**lighthouse_item)
                except (KeyError, ValueError):
                    continue

            try:
                for item in audits["layout-shift-elements"]['details']['items']:
                    node = item["node"]
                    lighthouse_item = {}
                    for aspect in node:
                        if aspect != "boundingRect":
                            if type(node[aspect]) == str:
                                if '{' in node[aspect] or '}' in node[aspect]:
                                    continue
                            lighthouse_item[aspect] = round(node[aspect], 2) if type(node[aspect]) == float \
                                else node[aspect]
                    lighthouse_item["report"] = lighthouse_report

                    LighthouseReportLayoutElementsItem.objects.create(**lighthouse_item)
            except (KeyError, ValueError):
                continue

    file_name = f'lighthouse_report_{domain}.xlsx'
    output = f'{settings.REPORT_PATH}/{file_name}'
    writer = pd.ExcelWriter(output, engine='xlsxwriter')
    keywords = df['Keyword']
    df = pd.DataFrame(
        {
            'Requested': [amount],
            'Requested from API': [len(keywords)],
        }
    )
    df.to_excel(writer, sheet_name='requested', index=False)

    df = pd.DataFrame(
        {
            'Total URLs': [len(list(set(keywords))) * 10],
            'Failed': [len(failed)],
            'Completed': [len(completed)]
        }
    )
    df.to_excel(writer, sheet_name='results', index=False)

    failed = pd.DataFrame({'Keyword': failed})
    failed.to_excel(writer, sheet_name='failed', index=False)
    completed = pd.DataFrame({'Keyword': completed})
    completed.to_excel(writer, sheet_name='completed', index=False)
    writer.save()

    send_mail('Lighthouse Run', f'Lighthouse Run for domain {domain} is completed. ', output)


def create_lighthouse_ppt(df, response, domain):
    domain = domain.lower()

    features = ['position', 'cumulative_ls', 'largest_cp', 'max_potential_fid']
    corr_df = df[features].corr()

    unique_content_types = list(set(df['content_type']))
    content_types = Counter(df['content_type'])

    prs = Presentation(f"{settings.REPORT_PATH}/presentation_template.pptx")

    slide = prs.slides.add_slide(prs.slide_layouts[8])
    shapes = slide.shapes

    shapes.title.text = f'Lighthouse Report'
    shapes.placeholders[1].text = f'DOMAIN - {domain}'

    slide = prs.slides.add_slide(prs.slide_layouts[13])
    shapes = slide.shapes
    shapes.title.text = 'Web Core Vitals'

    left = Inches(2.5)
    top = Inches(3.0)
    shapes.add_picture(f"{settings.REPORT_PATH}/web_core_vitals.png", left, top,
                       width=Inches(8), height=Inches(3.5))

    slide = prs.slides.add_slide(prs.slide_layouts[13])
    shapes = slide.shapes

    shapes.title.text = 'Correlation'

    rows = 4
    cols = 2
    left = Inches(2)
    top = Inches(3)
    width = Inches(10)
    height = Inches(2)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(6)

    table.cell(0, 0).text = ''
    table.cell(0, 1).text = 'Position'
    row_num = 1
    for index, value in zip(features[1:], corr_df['position'][1:]):
        table.cell(row_num, 0).text = str(index)
        table.cell(row_num, 1).text = str(round(value, 2))
        row_num += 1

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes

    shapes.title.text = 'Content Types'
    chart_data = ChartData()
    chart_data.categories = unique_content_types
    values = ()
    for content_type in unique_content_types:
        values += (content_types[content_type] / len(df['content_type']),)
    chart_data.add_series('Series 1', values)

    x, y, cx, cy = Inches(4), Inches(3), Inches(5), Inches(4)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    wasted_percent, wasted_byte = 0, 0
    images_with_wasted_byte = 0
    run_time, status_codes_200, transfer_size, resource_size = 0, 0, 0, 0
    lcp_num_data, fid_num_data, cls_num_data = {}, {}, {}
    comp_lcp_num_data, comp_fid_num_data, comp_cls_num_data = {}, {}, {}
    num_domain, num_comp = 0, 0
    layout = {}

    directories_dict = {}
    added_urls = []
    for index, (url, keyword, lcp, fid, cls) in enumerate(zip(df['url'], df['keyword'], df['largest_cp'],
                                                              df['max_potential_fid'], df['cumulative_ls'])):
        if url.find(domain) != -1 and url not in added_urls:
            directories = generate_directories(url)
            search_volume = int(get_keyword_volume(keyword))

            subdomain = get_subdomain(url)
            if subdomain and subdomain != 'www':
                directories.insert(0, subdomain)
            else:
                subdomain = None

            for dir_index, directory in enumerate(directories):
                if subdomain and dir_index == 0:
                    page_type = 'subdomain'
                else:
                    page_type = 'directory' if dir_index == 0 else 'subdirectory'

                if directory not in directories_dict:
                    directories_dict[directory] = {'urls': [(index, search_volume)], 'count': 1,
                                                   'page_type': page_type, 'fid_score': [fid],
                                                   'cls_score': [cls], 'lcp_score': [lcp]}
                else:
                    directories_dict[directory]['urls'].append((index, search_volume))
                    directories_dict[directory]['count'] += 1
                    directories_dict[directory]['fid_score'].append(fid)
                    directories_dict[directory]['cls_score'].append(cls)
                    directories_dict[directory]['lcp_score'].append(lcp)

        added_urls.append(url)

    directories = [(key, directories_dict[key]['count']) for key in directories_dict]

    sorted_directories_dict = sorted(directories, key=itemgetter(1), reverse=True)
    sorted_indexes = []
    for key, _ in sorted_directories_dict:
        urls_sorted = [index for index, _ in sorted(directories_dict[key]['urls'], key=itemgetter(1), reverse=True)]
        sorted_indexes += urls_sorted

    unique_indexes = []
    for sorted_index in sorted_indexes:
        if sorted_index not in unique_indexes:
            unique_indexes.append(sorted_index)

    sorted_indexes = unique_indexes

    first_cp_data, first_cp_comp_data = {}, {}
    speed_index_data, speed_index_comp_data = {}, {}
    interactive_data, interactive_comp_data = {}, {}
    total_bt_data, total_bt_comp_data = {}, {}

    unmin_js_wb, unmin_js_wp, unmin_js_tb = 0, 0, 0
    unmin_css_wb, unmin_css_wp, unmin_css_tb = 0, 0, 0
    unused_css_wb, unused_css_wp, unused_css_tb = 0, 0, 0

    for domain, keyword, url, type, lcp_score, fid_score, cls_score, position, cumulative_ls, largest_cp,\
            max_fid, first_cp, speed_index, interactive, total_blocking_time, created_at in zip(
            df['domain'], df['keyword'], df['url'], df['type'], df['lcp_score'],
            df['fid_score'], df['cls_score'], df['position'], df['cumulative_ls'],
            df['largest_cp'], df['max_potential_fid'], df['first_cp'], df['speed_index'],
            df['interactive'], df['total_blocking_time'], df['created_at']
    ):
        domain = domain.lower()

        report = LighthouseReport.objects.filter(domain=domain, keyword=keyword, url=url, type=type,
                                                 created_at=created_at).first()

        if interactive is None:
            interactive = 0

        if total_blocking_time is None:
            total_blocking_time = 0

        wasted_percent_item = 0
        wasted_bytes_item = 0
        offscreen_items = LighthouseReportOffscreenImagesItem.objects.filter(report=report)
        for offscreen_item in offscreen_items:
            wasted_percent_item += offscreen_item.wastedPercent
            wasted_bytes_item += offscreen_item.wastedBytes
            if offscreen_item.wastedBytes > 0:
                images_with_wasted_byte += 1

        if len(offscreen_items) != 0:
            wasted_percent += wasted_percent_item / len(offscreen_items)
            wasted_byte += wasted_bytes_item / len(offscreen_items)

        unminified_js_elems = LighthouseReportUnminifiedJS.objects.filter(report=report)
        unmin_js_wb_item, unmin_js_tb_item, unmin_js_wp_item = 0, 0, 0
        for unminified_js_elem in unminified_js_elems:
            unmin_js_tb_item += unminified_js_elem.totalBytes
            unmin_js_wb_item += unminified_js_elem.wastedBytes
            unmin_js_wp_item += unminified_js_elem.wastedPercent

        if len(unminified_js_elems) != 0:
            unmin_js_tb += unmin_js_tb_item / len(unminified_js_elems)
            unmin_js_wb += unmin_js_wb_item / len(unminified_js_elems)
            unmin_js_wp += unmin_js_wp_item / len(unminified_js_elems)

        unminified_css_elems = LighthouseReportUnminifiedCSS.objects.filter(report=report)
        unmin_css_wb_item, unmin_css_tb_item, unmin_css_wp_item = 0, 0, 0
        for unminified_css_elem in unminified_css_elems:
            unmin_css_tb_item += unminified_css_elem.totalBytes
            unmin_css_wb_item += unminified_css_elem.wastedBytes
            unmin_css_wp_item += unminified_css_elem.wastedPercent

        if len(unminified_css_elems) != 0:
            unmin_css_tb += unmin_css_tb_item / len(unminified_css_elems)
            unmin_css_wb += unmin_css_wb_item / len(unminified_css_elems)
            unmin_css_wp += unmin_css_wp_item / len(unminified_css_elems)

        unused_css_elems = LighthouseReportUnusedCSSRule.objects.filter(report=report)
        unused_css_wb_item, unused_css_tb_item, unused_css_wp_item = 0, 0, 0
        for unused_css_elem in unused_css_elems:
            unused_css_tb_item += unused_css_elem.totalBytes
            unused_css_wb_item += unused_css_elem.wastedBytes
            unused_css_wp_item += unused_css_elem.wastedPercent

        if len(unused_css_elems) != 0:
            unused_css_tb += unused_css_tb_item / len(unused_css_elems)
            unused_css_wb += unused_css_wb_item / len(unused_css_elems)
            unused_css_wp += unused_css_wp_item / len(unused_css_elems)

        network_items = LighthouseReportNetworkReqItems.objects.filter(report=report)
        run_time_item, status_code_200_item, transfer_size_item, resource_size_item = 0, 0, 0, 0
        for network_item in network_items:
            if network_item.startTime is None:
                run_time_item += network_item.endTime
            elif network_item.endTime is None:
                continue
            else:
                run_time_item += network_item.endTime - network_item.startTime

            if network_item.statusCode == 200:
                status_code_200_item += 1
            transfer_size_item += network_item.transferSize
            resource_size_item += network_item.resourceSize

        if len(network_items) != 0:
            run_time += run_time_item / len(network_items)
            status_codes_200 += status_code_200_item / len(network_items)
            transfer_size += transfer_size_item / len(network_items)
            resource_size += resource_size_item / len(network_items)

        if url.find(domain) != -1:
            if 0 < position <= 10:
                if position not in cls_num_data:
                    cls_num_data[position] = cumulative_ls
                else:
                    cls_num_data[position] += cumulative_ls

                if position not in first_cp_data:
                    first_cp_data[position] = first_cp
                else:
                    first_cp_data[position] += first_cp

                if position not in speed_index_data:
                    speed_index_data[position] = speed_index
                else:
                    speed_index_data[position] += speed_index

                if position not in interactive_data:
                    interactive_data[position] = interactive
                else:
                    interactive_data[position] += interactive

                if position not in total_bt_data:
                    total_bt_data[position] = total_blocking_time
                else:
                    total_bt_data[position] += total_blocking_time

                if position not in fid_num_data:
                    fid_num_data[position] = max_fid
                else:
                    fid_num_data[position] += max_fid

                if position not in lcp_num_data:
                    lcp_num_data[position] = largest_cp
                else:
                    lcp_num_data[position] += largest_cp

            layout_elements = LighthouseReportLayoutElementsItem.objects.filter(report=report)

            directories = generate_directories(url)
            subdomain = get_subdomain(url)
            paths = [element.path for element in layout_elements]
            if subdomain and subdomain != 'www':
                if subdomain not in layout:
                    layout[subdomain] = {'page_type': 'subdomain', 'consistency': [len(layout_elements)],
                                         'paths': paths, 'layout_elements': [layout_elements], 'count': 1}
                else:
                    layout[subdomain]['consistency'].append(len(layout_elements))
                    layout[subdomain]['paths'] += paths
                    layout[subdomain]['layout_elements'].append(report.id)
                    layout[subdomain]['count'] += 1

            for index, directory in enumerate(directories):
                page_type = 'directory' if index == 0 else 'subdirectory'
                if directory not in layout:
                    layout[directory] = {'page_type': page_type, 'consistency': [len(layout_elements)],
                                         'paths': paths, 'layout_elements': [layout_elements], 'count': 1}
                else:
                    layout[directory]['consistency'].append(len(layout_elements))
                    layout[directory]['paths'] += paths
                    layout[directory]['layout_elements'].append(layout_elements)
                    layout[directory]['count'] += 1

            num_domain += 1
        else:
            if 0 < position <= 10:
                if position not in comp_cls_num_data:
                    comp_cls_num_data[position] = cumulative_ls
                else:
                    comp_cls_num_data[position] += cumulative_ls

                if position not in first_cp_comp_data:
                    first_cp_comp_data[position] = first_cp
                else:
                    first_cp_comp_data[position] += first_cp

                if position not in speed_index_comp_data:
                    speed_index_comp_data[position] = speed_index
                else:
                    speed_index_comp_data[position] += speed_index

                if position not in interactive_comp_data:
                    interactive_comp_data[position] = interactive
                else:
                    interactive_comp_data[position] += interactive

                if position not in total_bt_comp_data:
                    total_bt_comp_data[position] = total_blocking_time
                else:
                    total_bt_comp_data[position] += total_blocking_time

                if position not in comp_fid_num_data:
                    comp_fid_num_data[position] = max_fid
                else:
                    comp_fid_num_data[position] += max_fid

                if position not in comp_lcp_num_data:
                    comp_lcp_num_data[position] = largest_cp
                else:
                    comp_lcp_num_data[position] += largest_cp
            num_comp += 1

    create_graph_slide(lcp_num_data, num_domain, num_comp, 'LCP', prs, comp_lcp_num_data, metric_type='Value')
    create_graph_slide(fid_num_data, num_domain, num_comp, 'Max FID', prs, comp_fid_num_data, metric_type='Value')
    create_graph_slide(cls_num_data, num_domain, num_comp, 'CLS', prs, comp_cls_num_data, metric_type='Value')
    create_graph_slide(first_cp_data, num_domain, num_comp, 'First Contentful Paint',
                       prs, first_cp_comp_data, metric_type='Value')
    create_graph_slide(speed_index_data, num_domain, num_comp, 'Speed Index', prs,
                       speed_index_comp_data, metric_type='Value')
    create_graph_slide(interactive_data, num_domain, num_comp, 'Time To Interactive',
                       prs, interactive_comp_data, metric_type='Value')
    create_graph_slide(total_bt_data, num_domain, num_comp, 'Total Blocking Time',
                       prs, total_bt_comp_data, metric_type='Value')

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes
    shapes.title.text = 'Directory Lighthouse Scores'

    left = Inches(0.5)
    top = Inches(3)
    width = Inches(10)
    height = Inches(1)
    rows = 12
    cols = 6

    table = shapes.add_table(rows, cols, left, top, width, height).table

    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(1)
    table.columns[3].width = Inches(2)
    table.columns[4].width = Inches(2)
    table.columns[5].width = Inches(2)

    table.cell(0, 0).text = 'Element'
    table.cell(0, 1).text = 'Page Type'
    table.cell(0, 2).text = 'Number of occurences'
    table.cell(0, 3).text = 'LCP Value'
    table.cell(0, 4).text = 'Max FID Value'
    table.cell(0, 5).text = 'CLS Value'

    for index, (key, _) in enumerate(sorted_directories_dict[:11], 1):
        table.cell(index, 0).text = key if key else '/'
        table.cell(index, 1).text = directories_dict[key]['page_type']
        table.cell(index, 2).text = str(directories_dict[key]['count'])
        table.cell(index, 3).text = str(average(directories_dict[key]['lcp_score']))
        table.cell(index, 4).text = str(average(directories_dict[key]['fid_score']))
        table.cell(index, 5).text = str(average(directories_dict[key]['cls_score']))

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes
    shapes.title.text = 'CLS Consistency'

    rows = 11
    cols = 3

    table = shapes.add_table(rows, cols, left, top, width, height).table

    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(3)
    table.columns[1].width = Inches(3)

    table.cell(0, 0).text = 'Page Type'
    table.cell(0, 1).text = 'Element'
    table.cell(0, 2).text = 'Consistency Score'

    consistency_keys = []
    for key in layout:
        common_paths = sorted(
            set(layout[key]['paths']),
            key=lambda ele: layout[key]['paths'].count(ele),
            reverse=True
        )[:int(average(layout[key]['consistency']))]
        consistency_score_overall = 0
        for layout_elements in layout[key]['layout_elements']:
            consistency_score = 0
            try:
                for layout_element in layout_elements:
                    if layout_element.path in common_paths:
                        consistency_score += 1
            except TypeError:
                continue

            consistency_score = consistency_score / len(common_paths) if len(common_paths) != 0 else 0
            consistency_score_overall += consistency_score
        consistency_score_overall = round(consistency_score_overall / len(layout[key]['layout_elements']), 2) \
            if len(layout[key]['layout_elements']) != 0 else 0
        consistency_keys.append((key, consistency_score_overall, layout[key]['count']))

    consistency_sorted = sorted(consistency_keys, key=itemgetter(2), reverse=True)

    for index, (key, consistency, _) in enumerate(consistency_sorted[:9], 1):
        table.cell(index, 0).text = layout[key]['page_type']
        if not key:
            key = '/'
        table.cell(index, 1).text = key
        table.cell(index, 2).text = str(consistency)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes
    shapes.title.text = 'Timelines Order'

    left = Inches(0.5)
    top = Inches(3)
    width = Inches(10)
    height = Inches(1)
    rows = 8
    cols = 4

    table = shapes.add_table(rows, cols, left, top, width, height).table

    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(1)
    table.columns[2].width = Inches(1)
    table.columns[3].width = Inches(8)

    table.cell(0, 0).text = 'Element'
    table.cell(0, 1).text = 'Page Type'
    table.cell(0, 2).text = 'Number of occurences'
    table.cell(0, 3).text = 'Urls'

    for index, (key, _) in enumerate(sorted_directories_dict[:7], 1):
        table.cell(index, 0).text = key if key else '/'
        table.cell(index, 1).text = directories_dict[key]['page_type']
        table.cell(index, 2).text = str(directories_dict[key]['count'])
        urls_sorted = [index for index, _ in sorted(directories_dict[key]['urls'], key=itemgetter(1), reverse=True)]
        urls = [df['url'].iloc[url] for url in urls_sorted][:5]
        table.cell(index, 3).text = ", ".join(urls)

    for index in sorted_indexes:
        domain, keyword, url, type = df['domain'].iloc[index], df['keyword'].iloc[index], df['url'].iloc[index], \
            df['type'].iloc[index]
        largest_cp, max_fid, cumulative_ls = df['largest_cp'].iloc[index], df['max_potential_fid'].iloc[index], \
            df['cumulative_ls'].iloc[index]
        total_bt, server_rt, position = df['total_blocking_time'].iloc[index], df['server_response_time'].iloc[index], \
            df['position'].iloc[index]
        report = LighthouseReport.objects.filter(domain=domain, keyword=keyword, url=url, type=type).first()
        items = LighthouseReportItem.objects.filter(report=report)

        slide = prs.slides.add_slide(prs.slide_layouts[0])
        shapes = slide.shapes

        shapes.title.text = f'Keyword - {keyword}. Mobile Timeline'

        inch = 1
        text_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(4),
            width=Inches(8),
            height=Inches(2)
        )
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.font.size = Pt(15)

        left = Inches(inch)
        top = Inches(4.5)
        shapes.add_picture(f"{settings.REPORT_PATH}/timeline.PNG", left, top, width=Inches(10.8), height=Inches(0.25))

        largest_cp_placed, max_fid_placed = False, False
        largest_cp_placed_pos = 0

        images = []
        for item in items:
            image_path = f"{settings.REPORT_PATH}/{generate_random_number()}.jpg"
            images.append(image_path)
            with open(image_path, 'wb') as f:
                data = f'{item.data.replace("base64,", "")}'
                image = base64.b64decode(data)
                f.write(image)

        images = [Image.open(image) for image in images]
        new_im = Image.new('RGB', (1385, 225), color=(255, 255, 255))

        x_offset = 0
        for im in images:
            new_im.paste(im, (x_offset, 0))
            x_offset += im.size[0] + 20

        image_path = f'{settings.REPORT_PATH}/{domain}_{keyword}_{type}_{generate_random_number()}.jpg'
        new_im.save(image_path)
        shapes.add_picture(image_path, left, top=Inches(3.2), width=Inches(10.8), height=Inches(1.3))

        for item in items:
            text_box = slide.shapes.add_textbox(
                left=Inches(inch + 0.5),
                top=Inches(4.5),
                width=Inches(0.9),
                height=Inches(2)
            )
            tf = text_box.text_frame
            p = tf.add_paragraph()
            p.text = str(item.timing)
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0, 0, 0)

            if not largest_cp_placed:
                if largest_cp <= item.timing:
                    left = Inches(inch + 0.3)
                    top = Inches(4.8)
                    shapes.add_picture(f"{settings.REPORT_PATH}/timeline-pointer.PNG", left, top, width=Inches(0.3),
                                       height=Inches(0.3))

                    place_largest_cp(slide, inch, largest_cp)
                    largest_cp_placed_pos = inch
                    largest_cp_placed = True

            if not max_fid_placed:
                if max_fid < item.timing:
                    left = Inches(inch + 0.3)
                    top = Inches(4.8)
                    shapes.add_picture(f"{settings.REPORT_PATH}/timeline-pointer.PNG", left, top, width=Inches(0.3),
                                       height=Inches(0.3))

                    place_max_fid(slide, inch, max_fid, largest_cp_placed_pos)
                    max_fid_placed = True

            inch += 1.1

        if not largest_cp_placed:
            left = Inches(inch - 0.8)
            top = Inches(4.8)
            shapes.add_picture(f"{settings.REPORT_PATH}/timeline-pointer.PNG", left, top, width=Inches(0.3),
                               height=Inches(0.3))

            place_largest_cp(slide, inch - 1.1, largest_cp)

        if not max_fid_placed:
            left = Inches(inch - 0.8)
            top = Inches(4.8)
            shapes.add_picture(f"{settings.REPORT_PATH}/timeline-pointer.PNG", left, top, width=Inches(0.3),
                               height=Inches(0.3))

            place_max_fid(slide, inch - 1.1, max_fid, largest_cp_placed_pos)

        text_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(5.7),
            width=Inches(3),
            height=Inches(1)
        )
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.font.size = Pt(15)
        p.text = f'Total Blocking Time - {total_bt}'
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        text_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(5.9),
            width=Inches(3),
            height=Inches(1)
        )
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.font.size = Pt(15)
        p.text = f'Cumulative Layout Shift - {cumulative_ls}'
        p.font.bold = True
        if cumulative_ls <= 0.1:
            color = RGBColor(0, 255, 0)
        elif cumulative_ls <= 0.25:
            color = RGBColor(255, 255, 0)
        else:
            color = RGBColor(255, 0, 0)
        p.font.color.rgb = color

        text_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(6.1),
            width=Inches(3),
            height=Inches(1)
        )
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.font.size = Pt(15)
        p.text = f'Server Response Time - {server_rt}'
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        text_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(6.5),
            width=Inches(10.8),
            height=Inches(1)
        )
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.text = f'{url}'
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes

    shapes.title.text = f'Lighthouse Image Stats'

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = f'Average Wasted Percent: {round(wasted_percent / len(df["url"]), 2)}'

    p = tf.add_paragraph()
    p.text = f'Average Wasted Bytes: {round(wasted_byte / len(df["url"]), 2)}'

    p = tf.add_paragraph()
    p.text = f'Average Of Images With Wasted Bytes: {round(images_with_wasted_byte / len(df["url"]), 2)}'

    create_slide_for_static_stats(prs, df, unmin_js_wp, unmin_js_wb, unmin_js_tb, title="Unminified JS Stats")
    create_slide_for_static_stats(prs, df, unmin_css_wp, unmin_css_wb, unmin_css_tb, title="Unminified CSS Stats")
    create_slide_for_static_stats(prs, df, unused_css_wp, unused_css_wb, unused_css_tb, title="Unused CSS Stats")

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes
    shapes.title.text = 'Network Requests Stats'

    rows = 2
    cols = 4
    left = Inches(0.5)
    top = Inches(3)
    width = Inches(10)
    height = Inches(1)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(3)
    table.columns[3].width = Inches(3)

    table.cell(0, 0).text = 'Average Runtime'
    table.cell(0, 1).text = '% of 200 Status Codes'
    table.cell(0, 2).text = 'Average Transfer Size'
    table.cell(0, 3).text = 'Average Resource Size'

    table.cell(1, 0).text = str(round(run_time / len(df["url"]), 2))
    table.cell(1, 1).text = str(round(status_codes_200 / len(df["url"]), 2))
    table.cell(1, 2).text = str(round(transfer_size / len(df["url"]), 2))
    table.cell(1, 3).text = str(round(resource_size / len(df["url"]), 2))

    return prs.save(response)


def place_max_fid(slide, inch, max_fid, largest_cp_placed_pos):
    top = 4.9
    if largest_cp_placed_pos == inch:
        top += 0.5

    text_box = slide.shapes.add_textbox(
        left=Inches(inch),
        top=Inches(top),
        width=Inches(0.9),
        height=Inches(4)
    )

    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = f'Max FID\n{max_fid}'
    p.font.bold = True
    p.font.size = Pt(10)

    if max_fid <= 100:
        color = RGBColor(0, 255, 0)
    elif max_fid <= 300:
        color = RGBColor(255, 255, 0)
    else:
        color = RGBColor(255, 0, 0)
    p.font.color.rgb = color


def place_largest_cp(slide, inch, largest_cp):
    text_box = slide.shapes.add_textbox(
        left=Inches(inch),
        top=Inches(4.9),
        width=Inches(0.9),
        height=Inches(4)
    )

    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = f'LCP\n{largest_cp}'
    p.font.bold = True
    p.font.size = Pt(10)

    if largest_cp / 1000 <= 2.5:
        color = RGBColor(0, 255, 0)
    elif largest_cp / 1000 <= 4.0:
        color = RGBColor(255, 255, 0)
    else:
        color = RGBColor(255, 0, 0)

    p.font.color.rgb = color


def create_graph_slide(data, num_domain, num_comp, feature_name, prs, comp_data, metric_type='Score'):
    feature_data, rankings = [], []
    comp_feature_data, comp_rankings = [], []

    for i in range(1, 11):
        try:
            feature_data.append(round(data[i] / num_domain, 2))
            rankings.append(i)
        except KeyError:
            logger.info(f'PPT: No data for position {i}')

    for i in range(1, 11):
        try:
            comp_feature_data.append(round(comp_data[i] / num_comp, 2))
            comp_rankings.append(i)
        except KeyError:
            logger.info(f'PPT: No comp data for position {i}')

    fig = plt.figure()
    plt.axes(frame_on=False).yaxis.grid()
    plt.plot(rankings, feature_data, color='green', linewidth=2, label='Domain')
    fig.suptitle(f'{feature_name} {metric_type}s for Rankings', fontsize=18)
    plt.xlabel('Organic Rankings', fontsize=8)
    plt.ylabel(f'{feature_name} {metric_type}', fontsize=8)
    plt.bar(comp_rankings, comp_feature_data, color='grey', label='Competitors')
    _, y_top = plt.ylim()
    plt.ylim([0, y_top*1.15])
    plt.legend()
    filepath = f'{settings.REPORT_PATH}/{feature_name}_{generate_random_number()}.jpg'
    plt.savefig(filepath)

    slide = prs.slides.add_slide(prs.slide_layouts[30])
    shapes = slide.shapes

    left = Inches(2.5)
    top = Inches(1.0)
    shapes.add_picture(filepath, left, top, width=Inches(8), height=Inches(5.5))


def create_slide_for_static_stats(prs, df, wasted_percent, wasted_byte, total_bytes, title):
    if wasted_percent and wasted_byte and total_bytes:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        shapes = slide.shapes

        shapes.title.text = title

        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = f'Average Wasted Percent: {round(wasted_percent / len(df["url"]), 2)}'

        p = tf.add_paragraph()
        p.text = f'Average Wasted Bytes: {round(wasted_byte / len(df["url"]), 2)}'

        p = tf.add_paragraph()
        p.text = f'Average Total Bytes: {round(total_bytes / len(df["url"]), 2)}'

